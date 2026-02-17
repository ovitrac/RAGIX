"""
RAGIX Project RAG - File Ingestion and Normalization

Crawls project folders, detects file types, and converts heterogeneous
files (code, docs, Office, PDF) into normalized text for indexing.

Supported formats:
    - Code: .java, .py, .js, .ts, .c, .cpp, .go, .rs, etc.
    - Config: .xml, .json, .yaml, .properties, etc.
    - Docs: .md, .txt, .rst
    - Office: .docx, .pptx, .xlsx, .odt, .odp, .ods (via pandoc/python libs)
    - PDF: .pdf (via pdftotext)

Author: Olivier Vitrac, PhD, HDR | olivier.vitrac@adservio.fr | Adservio | 2025-12-09
"""

import fnmatch
import os
import re
import shutil
import subprocess
import tempfile
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Dict, Iterator, List, Optional, Set, Tuple, Any
import logging

from .config import RAGConfig, IndexingFilters, get_rag_dir
from .metadata import FileKind, FileMetadata

logger = logging.getLogger(__name__)

# =============================================================================
# File Type Detection
# =============================================================================

# Extension to FileKind mapping
EXTENSION_MAP: Dict[str, FileKind] = {
    # Code
    ".java": FileKind.CODE_JAVA,
    ".py": FileKind.CODE_PYTHON,
    ".pyw": FileKind.CODE_PYTHON,
    ".pyx": FileKind.CODE_PYTHON,
    ".pyi": FileKind.CODE_PYTHON,
    ".js": FileKind.CODE_JAVASCRIPT,
    ".mjs": FileKind.CODE_JAVASCRIPT,
    ".cjs": FileKind.CODE_JAVASCRIPT,
    ".jsx": FileKind.CODE_JAVASCRIPT,
    ".ts": FileKind.CODE_TYPESCRIPT,
    ".tsx": FileKind.CODE_TYPESCRIPT,
    ".c": FileKind.CODE_C,
    ".h": FileKind.CODE_C,
    ".cpp": FileKind.CODE_CPP,
    ".hpp": FileKind.CODE_CPP,
    ".cc": FileKind.CODE_CPP,
    ".cxx": FileKind.CODE_CPP,
    ".go": FileKind.CODE_GO,
    ".rs": FileKind.CODE_RUST,
    ".rb": FileKind.CODE_RUBY,
    ".php": FileKind.CODE_PHP,
    ".kt": FileKind.CODE_KOTLIN,
    ".kts": FileKind.CODE_KOTLIN,
    ".scala": FileKind.CODE_SCALA,
    ".swift": FileKind.CODE_SWIFT,
    ".sh": FileKind.CODE_SHELL,
    ".bash": FileKind.CODE_SHELL,
    ".zsh": FileKind.CODE_SHELL,
    ".fish": FileKind.CODE_SHELL,
    ".ps1": FileKind.CODE_SHELL,
    ".bat": FileKind.CODE_SHELL,
    ".cmd": FileKind.CODE_SHELL,
    ".sql": FileKind.CODE_SQL,
    ".ddl": FileKind.CODE_SQL,

    # Config
    ".xml": FileKind.CONFIG_XML,
    ".json": FileKind.CONFIG_JSON,
    ".yaml": FileKind.CONFIG_YAML,
    ".yml": FileKind.CONFIG_YAML,
    ".toml": FileKind.CONFIG_OTHER,
    ".ini": FileKind.CONFIG_OTHER,
    ".cfg": FileKind.CONFIG_OTHER,
    ".conf": FileKind.CONFIG_OTHER,
    ".properties": FileKind.CONFIG_PROPERTIES,
    ".env": FileKind.CONFIG_OTHER,

    # Build
    ".gradle": FileKind.BUILD_GRADLE,

    # Documentation
    ".md": FileKind.DOC_MARKDOWN,
    ".markdown": FileKind.DOC_MARKDOWN,
    ".txt": FileKind.DOC_TEXT,
    ".rst": FileKind.DOC_RST,
    ".pdf": FileKind.DOC_PDF,
    ".docx": FileKind.DOC_DOCX,
    ".doc": FileKind.DOC_DOCX,
    ".odt": FileKind.DOC_ODT,
    ".pptx": FileKind.DOC_PPTX,
    ".ppt": FileKind.DOC_PPTX,
    ".odp": FileKind.DOC_PPTX,
    ".xlsx": FileKind.DOC_XLSX,
    ".xls": FileKind.DOC_XLSX,
    ".ods": FileKind.DOC_XLSX,

    # Data
    ".csv": FileKind.DATA_CSV,
    ".xsd": FileKind.DATA_XSD,
}

# Special filename patterns
FILENAME_PATTERNS: Dict[str, FileKind] = {
    "pom.xml": FileKind.BUILD_POM,
    "build.gradle": FileKind.BUILD_GRADLE,
    "build.gradle.kts": FileKind.BUILD_GRADLE,
    "settings.gradle": FileKind.BUILD_GRADLE,
    "Makefile": FileKind.BUILD_MAKEFILE,
    "makefile": FileKind.BUILD_MAKEFILE,
    "CMakeLists.txt": FileKind.BUILD_CMAKE,
}


def detect_file_kind(file_path: Path) -> FileKind:
    """
    Detect file kind based on extension and filename.

    Args:
        file_path: Path to file

    Returns:
        FileKind enum value
    """
    filename = file_path.name.lower()
    ext = file_path.suffix.lower()

    # Check special filenames first
    for pattern, kind in FILENAME_PATTERNS.items():
        if filename == pattern.lower():
            return kind

    # Check extension
    if ext in EXTENSION_MAP:
        return EXTENSION_MAP[ext]

    return FileKind.UNKNOWN


def detect_language(file_path: Path, kind: FileKind) -> Optional[str]:
    """
    Detect programming language or document language.

    Args:
        file_path: Path to file
        kind: Detected file kind

    Returns:
        Language string (e.g., "java", "python", "en", "fr")
    """
    kind_str = kind.value

    if kind_str.startswith("code_"):
        return kind_str.replace("code_", "")

    # For docs, could detect language from content (future)
    return None


def is_code_file(kind: FileKind) -> bool:
    """Check if file kind is source code."""
    return kind.value.startswith("code_")


def is_doc_file(kind: FileKind) -> bool:
    """Check if file kind is documentation."""
    return kind.value.startswith("doc_")


def needs_conversion(kind: FileKind) -> bool:
    """Check if file needs conversion to text."""
    return kind in (
        FileKind.DOC_PDF,
        FileKind.DOC_DOCX,
        FileKind.DOC_PPTX,
        FileKind.DOC_XLSX,
        FileKind.DOC_ODT,
    )


# =============================================================================
# File Conversion
# =============================================================================

def convert_pdf_to_text(file_path: Path, timeout: int = 60) -> str:
    """
    Convert PDF to text using pdftotext.

    Args:
        file_path: Path to PDF file
        timeout: Conversion timeout in seconds

    Returns:
        Extracted text content
    """
    pdftotext_path = shutil.which("pdftotext")
    if not pdftotext_path:
        raise RuntimeError("pdftotext not found. Install with: sudo apt install poppler-utils")

    cmd = [pdftotext_path, "-layout", str(file_path), "-"]

    result = subprocess.run(
        cmd,
        capture_output=True,
        text=True,
        timeout=timeout,
    )

    if result.returncode != 0:
        raise RuntimeError(f"PDF conversion failed: {result.stderr}")

    return result.stdout


def convert_office_to_text(file_path: Path, timeout: int = 60) -> str:
    """
    Convert Office document to markdown using pandoc.

    Args:
        file_path: Path to Office file (.docx, .odt, etc.)
        timeout: Conversion timeout

    Returns:
        Extracted text as markdown
    """
    pandoc_path = shutil.which("pandoc")
    if not pandoc_path:
        raise RuntimeError("pandoc not found. Install with: sudo apt install pandoc")

    cmd = [pandoc_path, "-t", "markdown", str(file_path)]

    result = subprocess.run(
        cmd,
        capture_output=True,
        text=True,
        timeout=timeout,
    )

    if result.returncode != 0:
        raise RuntimeError(f"Document conversion failed: {result.stderr}")

    return result.stdout


def convert_pptx_to_text(file_path: Path) -> str:
    """
    Convert PowerPoint to text using python-pptx.

    Args:
        file_path: Path to PPTX file

    Returns:
        Extracted text with slide markers
    """
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("python-pptx not installed. Install with: pip install python-pptx")

    prs = Presentation(str(file_path))
    text_parts = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = [f"\n## Slide {slide_num}\n"]

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())

            if hasattr(shape, "has_table") and shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    slide_text.append(row_text)

        if len(slide_text) > 1:
            text_parts.append("\n".join(slide_text))

    return "\n\n".join(text_parts)


def convert_xlsx_to_text(file_path: Path) -> str:
    """
    Convert Excel to text using openpyxl.

    Args:
        file_path: Path to XLSX file

    Returns:
        Extracted text with sheet markers
    """
    try:
        from openpyxl import load_workbook
    except ImportError:
        raise RuntimeError("openpyxl not installed. Install with: pip install openpyxl")

    wb = load_workbook(str(file_path), read_only=True, data_only=True)
    text_parts = []

    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        sheet_text = [f"\n## Sheet: {sheet_name}\n"]

        for row in sheet.iter_rows(values_only=True):
            row_values = [str(cell) if cell is not None else "" for cell in row]
            if any(v.strip() for v in row_values):
                sheet_text.append(" | ".join(row_values))

        if len(sheet_text) > 1:
            text_parts.append("\n".join(sheet_text))

    wb.close()
    return "\n\n".join(text_parts)


def normalize_file(file_path: Path, kind: FileKind) -> str:
    """
    Read and normalize file content to text.

    Args:
        file_path: Path to file
        kind: Detected file kind

    Returns:
        Normalized text content

    Raises:
        RuntimeError: If conversion fails
    """
    # Handle files needing conversion
    if kind == FileKind.DOC_PDF:
        return convert_pdf_to_text(file_path)

    elif kind == FileKind.DOC_DOCX or kind == FileKind.DOC_ODT:
        return convert_office_to_text(file_path)

    elif kind == FileKind.DOC_PPTX:
        try:
            return convert_pptx_to_text(file_path)
        except RuntimeError:
            # Fallback to pandoc
            return convert_office_to_text(file_path)

    elif kind == FileKind.DOC_XLSX:
        try:
            return convert_xlsx_to_text(file_path)
        except RuntimeError:
            # xlsx doesn't have good pandoc support, return empty
            logger.warning(f"Could not convert {file_path}")
            return ""

    # Text-based files: read directly
    else:
        try:
            return file_path.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            # Try with latin-1 fallback
            try:
                return file_path.read_text(encoding="latin-1")
            except Exception as e:
                logger.warning(f"Could not read {file_path}: {e}")
                return ""


# =============================================================================
# Ingest Result
# =============================================================================

@dataclass
class IngestResult:
    """Result of ingesting a single file."""
    file_path: str
    success: bool
    file_metadata: Optional[FileMetadata] = None
    content: str = ""
    error: Optional[str] = None
    chunk_count: int = 0


# =============================================================================
# File Ingester
# =============================================================================

class FileIngester:
    """
    Crawls and ingests project files for RAG indexing.

    Handles:
        - File discovery with glob patterns
        - File type detection
        - Content normalization
        - Change detection
    """

    def __init__(
        self,
        project_root: Path,
        config: Optional[RAGConfig] = None,
    ):
        """
        Initialize file ingester.

        Args:
            project_root: Project root directory
            config: RAG configuration (loads from .RAG/config.yaml if None)
        """
        self.project_root = project_root
        self.config = config

        # Get filters from config
        if config:
            self.filters = config.indexing
        else:
            self.filters = IndexingFilters()

    def scan_files(self) -> Iterator[Path]:
        """
        Scan project for files matching filters.

        Yields:
            Path objects for matching files
        """
        for root, dirs, files in os.walk(self.project_root, followlinks=True):
            root_path = Path(root)
            rel_root = root_path.relative_to(self.project_root)

            # Filter directories (skip excluded)
            dirs[:] = [
                d for d in dirs
                if not self._is_excluded(rel_root / d)
            ]

            for filename in files:
                file_path = root_path / filename
                rel_path = file_path.relative_to(self.project_root)

                # Check exclusions
                if self._is_excluded(rel_path):
                    continue

                # Check inclusions
                if not self._is_included(rel_path):
                    continue

                # Check file size
                try:
                    if file_path.stat().st_size > self.filters.max_file_size_bytes:
                        logger.debug(f"Skipping large file: {rel_path}")
                        continue
                except OSError:
                    continue

                yield file_path

    def count_files(self) -> int:
        """Count total files to index (for progress estimation)."""
        return sum(1 for _ in self.scan_files())

    def ingest_file(
        self,
        file_path: Path,
        file_id: str,
        profile_name: str,
    ) -> IngestResult:
        """
        Ingest a single file.

        Args:
            file_path: Path to file
            file_id: Unique file ID
            profile_name: Name of indexing profile used

        Returns:
            IngestResult with file metadata and content
        """
        rel_path = str(file_path.relative_to(self.project_root))

        try:
            # Detect file type
            kind = detect_file_kind(file_path)
            language = detect_language(file_path, kind)

            # Get file stats
            stat = file_path.stat()
            file_hash = FileMetadata.compute_hash(file_path)
            mtime = datetime.fromtimestamp(stat.st_mtime).isoformat()

            # Normalize content
            content = normalize_file(file_path, kind)

            # Create file metadata
            file_meta = FileMetadata(
                file_id=file_id,
                path=rel_path,
                kind=kind.value,
                language=language,
                hash=file_hash,
                mtime=mtime,
                size_bytes=stat.st_size,
                profile=profile_name,
                indexed_at=datetime.now().isoformat(),
            )

            return IngestResult(
                file_path=rel_path,
                success=True,
                file_metadata=file_meta,
                content=content,
            )

        except Exception as e:
            logger.error(f"Failed to ingest {rel_path}: {e}")
            return IngestResult(
                file_path=rel_path,
                success=False,
                error=str(e),
            )

    def ingest_all(
        self,
        profile_name: str,
        file_id_generator,
    ) -> Iterator[IngestResult]:
        """
        Ingest all files in project.

        Args:
            profile_name: Indexing profile name
            file_id_generator: Callable that returns next file ID

        Yields:
            IngestResult for each file
        """
        for file_path in self.scan_files():
            file_id = file_id_generator()
            yield self.ingest_file(file_path, file_id, profile_name)

    def _is_excluded(self, rel_path: Path) -> bool:
        """Check if path matches exclusion patterns."""
        path_str = str(rel_path)

        for pattern in self.filters.exclude_globs:
            if fnmatch.fnmatch(path_str, pattern):
                return True
            # Also check each path component
            if fnmatch.fnmatch(rel_path.name, pattern.rstrip("/**")):
                return True

        return False

    def _is_included(self, rel_path: Path) -> bool:
        """Check if path matches inclusion patterns."""
        path_str = str(rel_path)

        for pattern in self.filters.include_globs:
            if fnmatch.fnmatch(path_str, pattern):
                return True

        return False

    def detect_changes(
        self,
        indexed_files: Dict[str, FileMetadata],
    ) -> Dict[str, List[str]]:
        """
        Detect file changes since last indexing.

        Args:
            indexed_files: Dict of path -> FileMetadata from last index

        Returns:
            Dict with 'new', 'modified', 'deleted' lists of paths
        """
        changes = {
            "new": [],
            "modified": [],
            "deleted": [],
        }

        indexed_paths = set(indexed_files.keys())
        current_paths = set()

        for file_path in self.scan_files():
            rel_path = str(file_path.relative_to(self.project_root))
            current_paths.add(rel_path)

            if rel_path not in indexed_paths:
                changes["new"].append(rel_path)
            else:
                # Check if modified
                old_meta = indexed_files[rel_path]
                try:
                    current_hash = FileMetadata.compute_hash(file_path)
                    if current_hash != old_meta.hash:
                        changes["modified"].append(rel_path)
                except Exception:
                    pass

        # Find deleted files
        changes["deleted"] = list(indexed_paths - current_paths)

        return changes


# =============================================================================
# Tag Extraction
# =============================================================================

def extract_tags(content: str, file_path: str, kind: FileKind) -> List[str]:
    """
    Extract tags from content for knowledge graph.

    Looks for patterns like:
        - SK02, SC04, SG01 (ACME-ERP service codes)
        - Class names, function names
        - Keywords from file path

    Args:
        content: File content
        file_path: Relative file path
        kind: File kind

    Returns:
        List of extracted tags
    """
    tags = set()

    # Extract service codes (SKxx, SCxx, SGxx pattern)
    service_codes = re.findall(r'\b(SK\d{2}|SC\d{2}|SG\d{2})\b', content, re.IGNORECASE)
    tags.update(code.upper() for code in service_codes)

    # Extract from file path
    path_parts = Path(file_path).parts
    for part in path_parts:
        # Skip common directories
        if part.lower() not in ('src', 'main', 'java', 'python', 'test', 'resources'):
            # CamelCase split
            words = re.findall(r'[A-Z][a-z]+|[a-z]+', part)
            if words:
                tags.add(part)

    # For code files, extract class/function names
    if is_code_file(kind):
        # Java/Kotlin class names
        class_names = re.findall(r'\bclass\s+(\w+)', content)
        tags.update(class_names)

        # Python class/function names
        py_defs = re.findall(r'\b(?:class|def)\s+(\w+)', content)
        tags.update(py_defs)

    return list(tags)[:20]  # Limit to 20 tags
