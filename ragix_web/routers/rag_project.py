"""
RAGIX Project RAG Router - Massive project-level RAG with ChromaDB

This router provides API endpoints for the Project RAG system (Level 1 RAG).
It handles persistent, project-wide indexing using ChromaDB vector store.

Two-Level RAG Architecture:
    - Level 1 (this router): Project RAG - massive, persistent, ChromaDB-based (.RAG/)
    - Level 2 (rag.py router): Chat RAG - light, session-scoped, BM25-based (.ragix/)

IMPORTANT: The project path is the path selected in the Dashboard "Project Selection",
NOT the sandbox/launch directory. The .RAG/ folder is created inside the selected project.

Author: Olivier Vitrac, PhD, HDR | olivier.vitrac@adservio.fr | Adservio | 2025-12-09
"""

import logging
from pathlib import Path
from typing import Dict, Any, Optional

import requests
from fastapi import APIRouter, HTTPException, Query
from pydantic import BaseModel

logger = logging.getLogger(__name__)

router = APIRouter(prefix="/api/rag/project", tags=["Project RAG"])

# Current project path (set from frontend when user selects a project)
_current_project_path: Optional[str] = None

# Cached RAGProject instances per project root
_project_cache: Dict[str, Any] = {}

# Project RAG injection enabled per session (default: True)
_project_rag_enabled: Dict[str, bool] = {}


def set_project_rag_store(storage_root: str = ""):
    """Set the storage root from server.py (not used for project path)."""
    # Note: storage_root is the sandbox, not the project to index
    pass


def set_current_project(project_path: str):
    """Set the current project path (called when user selects a project)."""
    global _current_project_path
    _current_project_path = project_path
    logger.info(f"Project RAG: current project set to {project_path}")


def get_current_project() -> Optional[str]:
    """Get the current project path."""
    return _current_project_path


def clear_project_cache(project_path: Optional[str] = None):
    """Clear cached RAGProject instance(s) to force fresh connections."""
    global _project_cache
    if project_path:
        key = str(Path(project_path).resolve())
        if key in _project_cache:
            try:
                _project_cache[key].close()
            except Exception:
                pass
            del _project_cache[key]
            logger.info(f"Cleared cache for {key}")
    else:
        # Clear all
        for project in _project_cache.values():
            try:
                project.close()
            except Exception:
                pass
        _project_cache.clear()
        logger.info("Cleared all project caches")


def _get_rag_project(project_path: Optional[str] = None, fresh: bool = False):
    """
    Get or create RAGProject instance (lazy loading).

    Args:
        project_path: Path to the project to index. If None, uses current project.
        fresh: If True, clear cache and create new instance.

    Returns None if rag_project module is not available or no project selected.
    """
    try:
        from ragix_core.rag_project import RAGProject
    except ImportError as e:
        logger.warning(f"rag_project module not available: {e}")
        return None

    # Use provided path or current project
    path = project_path or _current_project_path
    if not path:
        return None

    root = Path(path).resolve()
    if not root.exists():
        logger.warning(f"Project path does not exist: {root}")
        return None

    key = str(root)

    # Clear cache if fresh requested
    if fresh and key in _project_cache:
        clear_project_cache(project_path)

    if key not in _project_cache:
        _project_cache[key] = RAGProject(root)

    return _project_cache[key]


# =============================================================================
# Request/Response Models
# =============================================================================

class ProjectInitRequest(BaseModel):
    """Request to initialize project RAG."""
    project_path: Optional[str] = None  # Path to project (uses current if not provided)
    profile: str = "mixed_docs_code"
    project_name: Optional[str] = None


class ProjectIndexRequest(BaseModel):
    """Request to start indexing."""
    project_path: Optional[str] = None  # Path to project (uses current if not provided)
    full_reindex: bool = False


class ProjectQueryRequest(BaseModel):
    """Request to query the project index."""
    project_path: Optional[str] = None  # Path to project (uses current if not provided)
    query: str
    top_k: int = 10
    collection: str = "mixed"


# =============================================================================
# Endpoints
# =============================================================================

@router.post("/set-project")
async def set_project(project_path: str = Query(..., description="Path to the project to index")) -> Dict[str, Any]:
    """
    Set the current project path for RAG indexing.

    This should be called when the user selects a project in the Dashboard.
    The .RAG/ folder will be created inside this project path.
    """
    path = Path(project_path).resolve()
    if not path.exists():
        raise HTTPException(404, f"Project path does not exist: {project_path}")
    if not path.is_dir():
        raise HTTPException(400, f"Project path is not a directory: {project_path}")

    set_current_project(str(path))

    # Clear cache to get fresh instance with current collections
    clear_project_cache(str(path))

    # Check if RAG already exists for this project
    project = _get_rag_project(str(path))
    exists = project.exists() if project else False
    initialized = project.is_initialized() if project else False

    return {
        "status": "ok",
        "project_path": str(path),
        "rag_exists": exists,
        "rag_initialized": initialized,
    }


@router.post("/refresh")
async def refresh_project_cache(
    project_path: Optional[str] = Query(None, description="Project path (uses current if not provided)")
) -> Dict[str, Any]:
    """
    Refresh the project cache to pick up changes.

    Clears the cached RAGProject instance to force fresh ChromaDB connections.
    Useful after reindexing or when collections appear stale.
    """
    path = project_path or _current_project_path
    clear_project_cache(path)

    # Get fresh instance
    project = _get_rag_project(path)
    if not project:
        return {"status": "no_project", "message": "No project selected"}

    return {
        "status": "ok",
        "project_path": path,
        "initialized": project.is_initialized(),
    }


@router.get("/current")
async def get_current() -> Dict[str, Any]:
    """Get the current project path."""
    return {
        "project_path": _current_project_path,
    }


class ProjectRagToggleRequest(BaseModel):
    """Request to toggle Project RAG injection."""
    enabled: bool


@router.post("/toggle")
async def toggle_project_rag(
    request: ProjectRagToggleRequest,
    session_id: str = Query("default", description="Session ID"),
) -> Dict[str, Any]:
    """
    Toggle Project RAG context injection for chat.

    When enabled (default), Project RAG context from .RAG/ is automatically
    injected into chat messages. When disabled, only Chat RAG (.ragix/) is used.
    """
    _project_rag_enabled[session_id] = request.enabled
    logger.info(f"Project RAG injection {'enabled' if request.enabled else 'disabled'} for session {session_id}")

    return {
        "status": "ok",
        "session_id": session_id,
        "project_rag_enabled": request.enabled,
    }


def is_project_rag_enabled(session_id: str = "default") -> bool:
    """Check if Project RAG injection is enabled for a session."""
    return _project_rag_enabled.get(session_id, True)  # Default: enabled


@router.get("/status")
async def get_project_status(
    project_path: Optional[str] = Query(None, description="Project path (uses current if not provided)")
) -> Dict[str, Any]:
    """
    Get Project RAG status.

    Returns comprehensive status including:
        - exists: Whether .RAG/ directory exists
        - initialized: Whether index is ready
        - is_indexing: Whether background indexing is running
        - state: Indexing state (files, chunks, progress)
    """
    # Use provided path or current project
    path = project_path or _current_project_path

    if not path:
        return {
            "available": True,
            "project_path": None,
            "message": "No project selected. Use POST /set-project to select a project.",
        }

    project = _get_rag_project(path)
    if not project:
        return {
            "available": False,
            "project_path": path,
            "error": "Project RAG module not installed or project path invalid.",
        }

    try:
        status = project.get_status()
        status["available"] = True
        status["current_project"] = path
        return status
    except Exception as e:
        logger.error(f"Failed to get project status: {e}")
        return {
            "available": True,
            "project_path": path,
            "exists": project.exists(),
            "error": str(e),
        }


@router.post("/init")
async def initialize_project(request: ProjectInitRequest) -> Dict[str, Any]:
    """
    Initialize .RAG/ directory with configuration.

    Creates the .RAG/ folder and config.yaml with the specified profile.
    The .RAG/ is created INSIDE the selected project, not in the sandbox.

    Profiles:
        - docs_only: Focus on documentation files (md, txt, rst, etc.)
        - mixed_docs_code: Both docs and code (default)
        - code_only: Focus on source code files
    """
    # Use provided path or current project
    path = request.project_path or _current_project_path
    if not path:
        raise HTTPException(400, "No project selected. Use POST /set-project first.")

    project = _get_rag_project(path)
    if not project:
        raise HTTPException(503, "Project RAG module not available or invalid path")

    try:
        from ragix_core.rag_project import ProfileType

        # Map string to ProfileType enum
        profile_map = {
            "docs_only": ProfileType.DOCS_ONLY,
            "mixed_docs_code": ProfileType.MIXED_DOCS_CODE,
            "code_only": ProfileType.CODE_ONLY,
        }

        profile = profile_map.get(request.profile, ProfileType.MIXED_DOCS_CODE)

        config = project.initialize(
            profile=profile,
            project_name=request.project_name or Path(path).name,
        )

        return {
            "status": "initialized",
            "project_path": path,
            "rag_dir": str(project.rag_dir),
            "profile": request.profile,
            "project_name": config.project_name,
        }
    except Exception as e:
        logger.error(f"Failed to initialize project: {e}")
        raise HTTPException(500, f"Initialization failed: {e}")


@router.post("/index")
async def start_indexing(request: ProjectIndexRequest) -> Dict[str, Any]:
    """
    Start background indexing.

    This runs in a separate thread and doesn't block the API.
    Use GET /status to monitor progress.

    Args:
        project_path: Path to the project (uses current if not provided)
        full_reindex: If true, clears existing index and rebuilds from scratch
    """
    # Use provided path or current project
    path = request.project_path or _current_project_path
    if not path:
        raise HTTPException(400, "No project selected. Use POST /set-project first.")

    project = _get_rag_project(path)
    if not project:
        raise HTTPException(503, "Project RAG module not available or invalid path")

    if project.is_indexing():
        return {
            "status": "already_running",
            "project_path": path,
            "message": "Indexing is already in progress",
        }

    try:
        # Start indexing in background
        worker = project.start_indexing(full_reindex=request.full_reindex)

        return {
            "status": "started",
            "project_path": path,
            "full_reindex": request.full_reindex,
            "message": f"Background indexing started for {path}. Use GET /status to monitor progress.",
        }
    except Exception as e:
        logger.error(f"Failed to start indexing: {e}")
        raise HTTPException(500, f"Failed to start indexing: {e}")


@router.post("/stop")
async def stop_indexing() -> Dict[str, Any]:
    """Stop background indexing."""
    project = _get_rag_project()
    if not project:
        raise HTTPException(503, "Project RAG module not available")

    if not project.is_indexing():
        return {
            "status": "not_running",
            "message": "No indexing in progress",
        }

    try:
        success = project.stop_indexing(wait=False)
        return {
            "status": "stopping" if success else "failed",
            "message": "Indexing stop requested" if success else "Failed to stop indexing",
        }
    except Exception as e:
        logger.error(f"Failed to stop indexing: {e}")
        raise HTTPException(500, f"Failed to stop indexing: {e}")


@router.get("/progress")
async def get_indexing_progress() -> Dict[str, Any]:
    """
    Get current indexing progress.

    Returns detailed progress information during background indexing.
    """
    project = _get_rag_project()
    if not project:
        raise HTTPException(503, "Project RAG module not available")

    progress = project.get_indexing_progress()
    if not progress:
        return {
            "is_indexing": False,
            "message": "No indexing in progress or completed",
        }

    is_indexing = project.is_indexing()

    # If indexing just completed, clear cache to pick up new collections
    if not is_indexing and progress.status in ("completed", "error", "cancelled"):
        clear_project_cache(_current_project_path)

    return {
        "is_indexing": is_indexing,
        "status": progress.status,
        "files_total": progress.files_total,
        "files_processed": progress.files_processed,
        "chunks_indexed": progress.chunks_indexed,
        "current_file": progress.current_file,
        "progress_percent": progress.progress_percent,
        "elapsed_seconds": progress.elapsed_seconds,
        "error": progress.error,
    }


@router.post("/query")
async def query_project(request: ProjectQueryRequest) -> Dict[str, Any]:
    """
    Query the project RAG index.

    Returns relevant chunks with citations from the indexed project.
    """
    path = request.project_path or _current_project_path
    project = _get_rag_project(path)
    if not project:
        raise HTTPException(503, "Project RAG module not available or no project selected")

    if not project.is_initialized():
        return {
            "query": request.query,
            "results": [],
            "error": "Project RAG not initialized. Use POST /init first.",
        }

    try:
        from ragix_core.rag_project import CollectionType

        collection_map = {
            "docs": CollectionType.DOCS,
            "code": CollectionType.CODE,
            "mixed": CollectionType.MIXED,
        }
        collection = collection_map.get(request.collection, CollectionType.MIXED)

        results = project.query(
            query_text=request.query,
            top_k=request.top_k,
            collection=collection,
        )

        return {
            "query": request.query,
            "collection": request.collection,
            "top_k": request.top_k,
            "results": [
                {
                    "chunk_id": r.chunk_id,
                    "score": r.score,
                    "file_path": r.file_path,
                    "line_start": r.line_start,
                    "line_end": r.line_end,
                    "content": r.content[:1000],  # Truncate for response
                    "citation": r.get_citation(),
                }
                for r in results
            ],
            "count": len(results),
        }
    except Exception as e:
        logger.error(f"Query failed: {e}")
        raise HTTPException(500, f"Query failed: {e}")


@router.get("/stats")
async def get_project_stats(
    project_path: Optional[str] = Query(None, description="Project path (uses current if not provided)")
) -> Dict[str, Any]:
    """
    Get detailed project RAG statistics.

    Returns statistics about indexed files, chunks, vector store, and graph.
    """
    path = project_path or _current_project_path
    project = _get_rag_project(path)
    if not project:
        raise HTTPException(503, "Project RAG module not available or no project selected")

    if not project.is_initialized():
        return {
            "initialized": False,
            "message": "Project RAG not initialized",
        }

    try:
        stats = project.get_stats()
        stats["initialized"] = True
        stats["project_path"] = path
        return stats
    except Exception as e:
        logger.error(f"Failed to get stats: {e}")
        raise HTTPException(500, f"Failed to get stats: {e}")


@router.get("/popular-concepts")
async def get_popular_concepts(
    project_path: Optional[str] = Query(None, description="Project path (uses current if not provided)"),
    top_k: int = Query(30, ge=1, le=500),
) -> Dict[str, Any]:
    """
    Get popular concepts for exploration.

    Returns the most frequently mentioned concepts across the indexed project.
    Useful for starting exploration of a new codebase.
    """
    path = project_path or _current_project_path
    project = _get_rag_project(path)
    if not project:
        raise HTTPException(503, "Project RAG module not available or no project selected")

    if not project.is_initialized():
        return {
            "concepts": [],
            "message": "Project RAG not initialized",
        }

    try:
        from ragix_core.rag_project.graph import NodeType, EdgeType

        graph = project.graph
        graph.load()

        # Get all concepts
        concepts = {n.id: n for n in graph._nodes.values() if n.type == NodeType.CONCEPT}

        # Count mentions per concept
        concept_counts = {}
        for edge in graph._edges:
            if edge.type == EdgeType.MENTIONS:
                target = edge.target
                concept_counts[target] = concept_counts.get(target, 0) + 1

        # Sort by popularity and filter out generic terms
        skip_terms = {'doc', 'Old', 'com', 'fr', 'org', 'xml', 'java', 'pom.xml', 'test'}
        popular = []
        for cid, count in sorted(concept_counts.items(), key=lambda x: x[1], reverse=True):
            if cid in concepts:
                label = concepts[cid].label
                if label not in skip_terms and len(label) > 2:
                    popular.append({"concept": label, "mentions": count})
                    if len(popular) >= top_k:
                        break

        return {
            "concepts": popular,
            "total_concepts": len(concepts),
        }
    except Exception as e:
        logger.error(f"Failed to get popular concepts: {e}")
        return {"concepts": [], "error": str(e)}


@router.get("/detect-changes")
async def detect_changes(
    project_path: Optional[str] = Query(None, description="Project path (uses current if not provided)"),
) -> Dict[str, Any]:
    """
    Detect file changes since last indexing.

    Compares current file tree against indexed metadata to find:
    - New files: Present on disk but not in index
    - Modified files: mtime or size changed
    - Deleted files: In index but not on disk

    Returns summary for UI notification display.
    """
    path = project_path or _current_project_path
    project = _get_rag_project(path)
    if not project:
        raise HTTPException(503, "Project RAG module not available or no project selected")

    if not project.is_initialized():
        return {
            "has_changes": False,
            "initialized": False,
            "message": "Project RAG not initialized - no changes to detect",
        }

    try:
        changes = project.detect_changes()
        return {
            "has_changes": changes["has_changes"],
            "initialized": True,
            "summary": changes["summary"],
            "new_count": len(changes["new"]),
            "modified_count": len(changes["modified"]),
            "deleted_count": len(changes["deleted"]),
            "total_changes": len(changes["new"]) + len(changes["modified"]) + len(changes["deleted"]),
            # Include first few file paths for preview
            "new_preview": changes["new"][:5],
            "modified_preview": changes["modified"][:5],
            "deleted_preview": changes["deleted"][:5],
        }
    except Exception as e:
        logger.error(f"Failed to detect changes: {e}")
        return {
            "has_changes": False,
            "error": str(e),
        }


@router.post("/search-concept")
async def search_concept(
    concept: str = Query(..., description="Concept to search for"),
    top_k: int = Query(20, ge=1, le=100),
    project_path: Optional[str] = Query(None, description="Project path (uses current if not provided)"),
) -> Dict[str, Any]:
    """
    Search for a concept across the project.

    Returns both vector search results and knowledge graph relationships.
    """
    path = project_path or _current_project_path
    project = _get_rag_project(path)
    if not project:
        raise HTTPException(503, "Project RAG module not available or no project selected")

    if not project.is_initialized():
        return {
            "concept": concept,
            "results": [],
            "error": "Project RAG not initialized",
        }

    try:
        return project.search_concept(concept, top_k=top_k)
    except Exception as e:
        logger.error(f"Concept search failed: {e}")
        raise HTTPException(500, f"Concept search failed: {e}")


@router.post("/summarize")
async def summarize_concept(
    concept: str = Query(..., description="Concept to summarize"),
    top_k: int = Query(10, ge=1, le=50),
    project_path: Optional[str] = Query(None, description="Project path (uses current if not provided)"),
    model: Optional[str] = Query(None, description="LLM model to use (uses default from settings if not provided)"),
) -> Dict[str, Any]:
    """
    Generate an LLM-powered summary of a concept with citations.

    Retrieves relevant chunks from the index, then uses the configured LLM
    (via Ollama) to generate a summary with source citations.
    """
    path = project_path or _current_project_path
    project = _get_rag_project(path)
    if not project:
        return {"error": "Project RAG module not available or no project selected"}

    if not project.is_initialized():
        return {"error": "Project RAG not initialized"}

    try:
        # Get relevant chunks
        results = project.query(concept, top_k=top_k)

        if not results:
            return {
                "concept": concept,
                "summary": f"No relevant content found for '{concept}' in the project index.",
                "citations": [],
            }

        # Build context for LLM
        context_parts = []
        citations = []
        for i, r in enumerate(results[:10]):  # Limit to 10 for prompt size
            citation = f"[{i+1}] {r.file_path}:{r.line_start}-{r.line_end}"
            context_parts.append(f"{citation}\n{r.content[:800]}")
            citations.append({
                "file_path": r.file_path,
                "line_start": r.line_start,
                "line_end": r.line_end,
                "score": r.score,
            })

        context_text = "\n\n---\n\n".join(context_parts)

        # Use OllamaLLM backend (same as chat tab)
        try:
            from ragix_core.llm_backends import OllamaLLM

            # Use provided model or default to a sensible model
            llm_model = model or "mistral"

            llm = OllamaLLM(model=llm_model)

            system_prompt = """You are a technical analyst specializing in software documentation and code analysis.
Your task is to provide concise, accurate summaries with proper source citations.
Always cite your sources using [N] notation matching the source numbers provided."""

            user_prompt = f"""Based on the following code/documentation excerpts about "{concept}", provide a concise summary (2-4 paragraphs) that:

1. Explains what "{concept}" is and its purpose in the project
2. Describes how it's used or implemented
3. Notes any important relationships or dependencies

## Sources:
{context_text}

## Summary:"""

            # Call LLM
            summary = llm.generate(
                system_prompt=system_prompt,
                history=[{"role": "user", "content": user_prompt}]
            )

            return {
                "concept": concept,
                "summary": summary,
                "citations": citations,
                "chunks_used": len(citations),
                "model": llm_model,
            }

        except ImportError as e:
            # LLM backend not available
            logger.error(f"LLM backend not available: {e}")
            return {
                "concept": concept,
                "summary": f"LLM backend not available. Found {len(results)} relevant chunks.",
                "citations": citations,
                "chunks_used": len(citations),
                "note": "Ensure ragix_core.llm_backends is available",
            }
        except requests.exceptions.ConnectionError:
            # Ollama not running
            return {
                "concept": concept,
                "summary": f"Ollama not running. Start it with 'ollama serve'. Found {len(results)} relevant chunks.",
                "citations": citations,
                "chunks_used": len(citations),
                "note": "Start Ollama: ollama serve",
            }
        except Exception as llm_err:
            logger.error(f"LLM call failed: {llm_err}")
            return {
                "concept": concept,
                "summary": f"LLM call failed: {str(llm_err)}. Found {len(results)} relevant chunks.",
                "citations": citations,
                "chunks_used": len(citations),
            }

    except Exception as e:
        logger.error(f"Summarization failed: {e}")
        return {"error": str(e)}


@router.delete("/clear")
async def clear_project_rag() -> Dict[str, Any]:
    """
    Clear all project RAG data.

    Removes all indexed data, metadata, and graph. The .RAG/ directory
    structure is preserved but emptied.
    """
    project = _get_rag_project()
    if not project:
        raise HTTPException(503, "Project RAG module not available")

    try:
        success = project.clear()
        return {
            "status": "cleared" if success else "failed",
            "message": "All project RAG data cleared" if success else "Failed to clear data",
        }
    except Exception as e:
        logger.error(f"Failed to clear project: {e}")
        raise HTTPException(500, f"Failed to clear: {e}")


# =============================================================================
# Concept Exploration Workflow (v0.33.1)
# =============================================================================

class ConceptExploreRequest(BaseModel):
    """Request for concept exploration."""
    query: str
    top_k: int = 20
    collection: str = "mixed"


@router.post("/concept-explore")
async def concept_explore(
    request: ConceptExploreRequest,
    project_path: Optional[str] = Query(None, description="Project path (uses current if not provided)"),
) -> Dict[str, Any]:
    """
    Explore a concept with dual view data (file-centric + graph-ready).

    Returns:
        - Files aggregated by relevance (chunk count, max score)
        - Recurring tags extracted from results
        - Suggested concept if not in graph
    """
    path = project_path or _current_project_path
    project = _get_rag_project(path)
    if not project:
        raise HTTPException(503, "Project RAG module not available or no project selected")

    if not project.is_initialized():
        return {"error": "Project RAG not initialized", "query": request.query}

    try:
        # Query vector store
        results = project.query(request.query, top_k=request.top_k, collection=request.collection)

        if not results:
            return {
                "query": request.query,
                "total_chunks": 0,
                "recurring_tags": [],
                "files": [],
                "suggested_concept": None,
            }

        # Aggregate by file
        files_map: Dict[str, Dict] = {}
        tag_counts: Dict[str, int] = {}

        for r in results:
            file_path = r.file_path
            if file_path not in files_map:
                files_map[file_path] = {
                    "path": file_path,
                    "chunks": [],
                    "max_score": 0.0,
                    "total_score": 0.0,
                }

            files_map[file_path]["chunks"].append({
                "chunk_id": r.chunk_id,
                "score": round(r.score, 4),
                "line_start": r.line_start,
                "line_end": r.line_end,
                "preview": r.content[:150] + "..." if len(r.content) > 150 else r.content,
                "kind": r.kind,
            })
            files_map[file_path]["max_score"] = max(files_map[file_path]["max_score"], r.score)
            files_map[file_path]["total_score"] += r.score

            # Extract tags from metadata
            tags_str = r.metadata.get("tags", "")
            if tags_str:
                for tag in tags_str.split(","):
                    tag = tag.strip()
                    if tag and len(tag) > 2:  # Skip very short tags
                        tag_counts[tag] = tag_counts.get(tag, 0) + 1

        # Sort files by max_score descending
        files_list = sorted(
            [
                {
                    **f,
                    "chunk_count": len(f["chunks"]),
                    "avg_score": round(f["total_score"] / len(f["chunks"]), 4),
                    "max_score": round(f["max_score"], 4),
                }
                for f in files_map.values()
            ],
            key=lambda x: (-x["max_score"], -x["chunk_count"])
        )

        # Get recurring tags (sorted by frequency)
        # Filter out generic terms
        generic_terms = {"doc", "old", "com", "fr", "org", "xml", "java", "pom", "test", "src", "main"}
        recurring_tags = [
            {"tag": tag, "count": count}
            for tag, count in sorted(tag_counts.items(), key=lambda x: -x[1])
            if tag.lower() not in generic_terms and count >= 2
        ][:20]  # Top 20 recurring tags

        # Check if query matches an existing concept
        suggested_concept = None
        try:
            graph = project.graph
            concept_node = graph.get_concept_by_label(request.query)
            if concept_node:
                suggested_concept = {
                    "label": request.query,
                    "exists_in_graph": True,
                    "node_id": concept_node.id,
                }
            else:
                # Suggest creating if query appears frequently in tags
                query_lower = request.query.lower()
                matching_tag = next((t for t in recurring_tags if t["tag"].lower() == query_lower), None)
                if matching_tag:
                    suggested_concept = {
                        "label": request.query,
                        "exists_in_graph": False,
                        "mentions": matching_tag["count"],
                    }
        except Exception as e:
            logger.debug(f"Graph lookup failed: {e}")

        return {
            "query": request.query,
            "total_chunks": len(results),
            "recurring_tags": recurring_tags,
            "files": files_list,
            "suggested_concept": suggested_concept,
        }

    except Exception as e:
        logger.error(f"Concept exploration failed: {e}")
        raise HTTPException(500, f"Concept exploration failed: {e}")


class ConceptGraphRequest(BaseModel):
    """Request for concept graph visualization."""
    concept: str
    depth: int = 2  # 1=chunks only, 2=chunks+files
    max_nodes: int = 50


@router.post("/concept-graph")
async def concept_graph(
    request: ConceptGraphRequest,
    project_path: Optional[str] = Query(None, description="Project path (uses current if not provided)"),
) -> Dict[str, Any]:
    """
    Get graph visualization data for a concept (D3.js force-directed layout).

    Returns nodes and edges for rendering a concept-centered graph:
    - Center: concept node
    - Ring 1: related chunks
    - Ring 2: parent files (if depth=2)
    """
    path = project_path or _current_project_path
    project = _get_rag_project(path)
    if not project:
        raise HTTPException(503, "Project RAG module not available or no project selected")

    if not project.is_initialized():
        return {"error": "Project RAG not initialized"}

    try:
        # Query for relevant chunks
        results = project.query(request.concept, top_k=request.max_nodes)

        nodes = []
        edges = []
        seen_nodes = set()

        # Add concept node at center
        concept_id = f"concept_{request.concept.lower().replace(' ', '_')}"
        nodes.append({
            "id": concept_id,
            "type": "concept",
            "label": request.concept,
            "group": 0,  # Center group
        })
        seen_nodes.add(concept_id)

        # Track files for second ring
        file_chunks: Dict[str, List[str]] = {}

        # Add chunk nodes (ring 1)
        for i, r in enumerate(results[:request.max_nodes]):
            chunk_id = f"chunk_{r.chunk_id[:8]}" if r.chunk_id else f"chunk_{i}"

            if chunk_id not in seen_nodes:
                nodes.append({
                    "id": chunk_id,
                    "type": "chunk",
                    "label": f"L{r.line_start}-{r.line_end}",
                    "file": r.file_path,
                    "score": round(r.score, 3),
                    "preview": r.content[:100] + "..." if len(r.content) > 100 else r.content,
                    "group": 1,  # Chunk ring
                    "line_start": r.line_start,
                    "line_end": r.line_end,
                })
                seen_nodes.add(chunk_id)

                # Edge from concept to chunk
                edges.append({
                    "source": concept_id,
                    "target": chunk_id,
                    "type": "mentions",
                    "value": r.score,  # Edge weight for layout
                })

                # Track file relationship
                if r.file_path not in file_chunks:
                    file_chunks[r.file_path] = []
                file_chunks[r.file_path].append(chunk_id)

        # Add file nodes (ring 2) if depth >= 2
        if request.depth >= 2:
            for file_path, chunk_ids in file_chunks.items():
                file_id = f"file_{hash(file_path) % 100000}"

                if file_id not in seen_nodes:
                    nodes.append({
                        "id": file_id,
                        "type": "file",
                        "label": Path(file_path).name,
                        "path": file_path,
                        "chunk_count": len(chunk_ids),
                        "group": 2,  # File ring
                    })
                    seen_nodes.add(file_id)

                # Edges from chunks to file
                for chunk_id in chunk_ids:
                    edges.append({
                        "source": chunk_id,
                        "target": file_id,
                        "type": "contains",
                        "value": 0.5,  # Weaker edge for layout
                    })

        # Find related concepts from the graph
        related_concepts = []
        try:
            graph = project.graph
            # Get concepts that appear in the same chunks
            tag_counts: Dict[str, int] = {}
            for r in results:
                tags_str = r.metadata.get("tags", "")
                if tags_str:
                    for tag in tags_str.split(","):
                        tag = tag.strip()
                        if tag and tag.lower() != request.concept.lower():
                            tag_counts[tag] = tag_counts.get(tag, 0) + 1

            related_concepts = [
                {"label": tag, "co_occurrences": count}
                for tag, count in sorted(tag_counts.items(), key=lambda x: -x[1])[:10]
            ]
        except Exception as e:
            logger.debug(f"Related concepts lookup failed: {e}")

        return {
            "concept": request.concept,
            "nodes": nodes,
            "edges": edges,
            "related_concepts": related_concepts,
            "stats": {
                "total_nodes": len(nodes),
                "total_edges": len(edges),
                "files": len(file_chunks),
                "chunks": len(results),
            },
        }

    except Exception as e:
        logger.error(f"Concept graph failed: {e}")
        raise HTTPException(500, f"Concept graph failed: {e}")


@router.get("/chunk/{chunk_id}/neighbors")
async def get_chunk_neighbors(
    chunk_id: str,
    project_path: Optional[str] = Query(None, description="Project path (uses current if not provided)"),
) -> Dict[str, Any]:
    """
    Get neighboring chunks for exploration.

    Returns:
        - Previous/next chunks in the same file
        - Semantically similar chunks (via vector similarity)
        - Concepts this chunk mentions
    """
    path = project_path or _current_project_path
    project = _get_rag_project(path)
    if not project:
        raise HTTPException(503, "Project RAG module not available or no project selected")

    if not project.is_initialized():
        return {"error": "Project RAG not initialized"}

    try:
        # Get the chunk content first (we need to find it)
        # Note: ChromaDB doesn't support direct ID lookup easily,
        # so we'll query by the chunk content or ID pattern

        # For now, return a placeholder - full implementation requires
        # metadata lookup by chunk_id
        return {
            "chunk_id": chunk_id,
            "neighbors": {
                "previous": None,
                "next": None,
                "similar": [],
            },
            "concepts": [],
            "note": "Full neighbor lookup requires chunk_id indexing (planned for v0.34)",
        }

    except Exception as e:
        logger.error(f"Chunk neighbors lookup failed: {e}")
        raise HTTPException(500, f"Chunk neighbors lookup failed: {e}")


# =================================================================
# Document conversion helpers for file viewer
# =================================================================

def _convert_document(file_path: Path) -> str:
    """Convert Word/ODT document to markdown text."""
    import shutil
    import subprocess

    # Try pandoc first
    pandoc_path = shutil.which("pandoc")
    if pandoc_path:
        try:
            result = subprocess.run(
                [pandoc_path, "-t", "markdown", str(file_path)],
                capture_output=True, text=True, timeout=60
            )
            if result.returncode == 0:
                return result.stdout
        except Exception as e:
            logger.warning(f"pandoc conversion failed: {e}")

    # Fallback: try python-docx for .docx
    if file_path.suffix.lower() in ['.docx']:
        try:
            from docx import Document
            doc = Document(str(file_path))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return '\n\n'.join(paragraphs)
        except ImportError:
            pass
        except Exception as e:
            logger.warning(f"python-docx failed: {e}")

    return f"[Document preview not available - install pandoc for full support]\n\nFile: {file_path.name}"


def _convert_presentation(file_path: Path) -> str:
    """Convert PowerPoint/ODP presentation to text."""
    suffix = file_path.suffix.lower()

    # Use python-pptx for .pptx
    if suffix in ['.pptx']:
        try:
            from pptx import Presentation
            prs = Presentation(str(file_path))
            text_parts = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"\n## Slide {slide_num}\n"]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                if len(slide_text) > 1:
                    text_parts.append("\n".join(slide_text))

            return "\n\n".join(text_parts) if text_parts else "[Empty presentation]"
        except ImportError:
            pass
        except Exception as e:
            logger.warning(f"python-pptx failed: {e}")

    # Try pandoc for other formats
    import shutil
    import subprocess
    pandoc_path = shutil.which("pandoc")
    if pandoc_path:
        try:
            result = subprocess.run(
                [pandoc_path, "-t", "markdown", str(file_path)],
                capture_output=True, text=True, timeout=60
            )
            if result.returncode == 0:
                return result.stdout
        except Exception as e:
            logger.warning(f"pandoc conversion failed: {e}")

    return f"[Presentation preview not available]\n\nFile: {file_path.name}"


def _convert_spreadsheet(file_path: Path) -> str:
    """Convert Excel/ODS spreadsheet to text table."""
    suffix = file_path.suffix.lower()

    # Use openpyxl for .xlsx
    if suffix in ['.xlsx']:
        try:
            from openpyxl import load_workbook
            wb = load_workbook(str(file_path), read_only=True, data_only=True)
            text_parts = []

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheet_text = [f"\n## Sheet: {sheet_name}\n"]

                for row in sheet.iter_rows(values_only=True, max_row=100):  # Limit rows
                    row_values = [str(cell) if cell is not None else "" for cell in row]
                    if any(v.strip() for v in row_values):
                        sheet_text.append(" | ".join(row_values))

                if len(sheet_text) > 1:
                    text_parts.append("\n".join(sheet_text))

            wb.close()
            return "\n\n".join(text_parts) if text_parts else "[Empty spreadsheet]"
        except ImportError:
            pass
        except Exception as e:
            logger.warning(f"openpyxl failed: {e}")

    # Try pandoc for other formats (ods)
    import shutil
    import subprocess
    pandoc_path = shutil.which("pandoc")
    if pandoc_path:
        try:
            result = subprocess.run(
                [pandoc_path, "-t", "markdown", str(file_path)],
                capture_output=True, text=True, timeout=60
            )
            if result.returncode == 0:
                return result.stdout
        except Exception as e:
            logger.warning(f"pandoc conversion failed: {e}")

    return f"[Spreadsheet preview not available]\n\nFile: {file_path.name}"


def _convert_pdf(file_path: Path) -> str:
    """Convert PDF to text using pdftotext."""
    import shutil
    import subprocess

    pdftotext_path = shutil.which("pdftotext")
    if pdftotext_path:
        try:
            result = subprocess.run(
                [pdftotext_path, "-layout", str(file_path), "-"],
                capture_output=True, text=True, timeout=120
            )
            if result.returncode == 0:
                return result.stdout
        except Exception as e:
            logger.warning(f"pdftotext failed: {e}")

    return f"[PDF preview not available - install poppler-utils for pdftotext]\n\nFile: {file_path.name}"


@router.get("/file-view")
async def get_file_view(
    file_path: str = Query(..., description="Path to the file (relative or absolute)"),
    project_path: Optional[str] = Query(None, description="Project path"),
    concept: Optional[str] = Query(None, description="Concept to highlight"),
) -> Dict[str, Any]:
    """
    Get file content with chunk information for the file viewer.

    Supports various formats:
    - Code files: displayed as-is with syntax highlighting
    - Documents: .docx, .odt converted via pandoc
    - Presentations: .pptx, .odp converted via python-pptx or pandoc
    - Spreadsheets: .xlsx, .ods converted via openpyxl or pandoc

    Returns:
        - File content (lines or converted text)
        - Chunks in this file with their line ranges
        - Concept associations for each chunk
        - Format type for appropriate rendering
    """
    import html as html_lib

    path = project_path or _current_project_path
    if not path:
        raise HTTPException(400, "No project path specified")

    # Resolve the file path
    project_root = Path(path)
    file_p = Path(file_path)

    # Handle both absolute and relative paths
    if file_p.is_absolute():
        actual_path = file_p
    else:
        actual_path = project_root / file_p

    # Security: ensure file is within project
    try:
        actual_path = actual_path.resolve()
        project_root = project_root.resolve()
        if not str(actual_path).startswith(str(project_root)):
            raise HTTPException(403, "File path outside project directory")
    except Exception:
        raise HTTPException(400, "Invalid file path")

    if not actual_path.exists():
        raise HTTPException(404, f"File not found: {file_path}")

    if not actual_path.is_file():
        raise HTTPException(400, "Path is not a file")

    # Determine file type and format
    suffix = actual_path.suffix.lower()

    # Document formats that need conversion
    doc_formats = {
        '.docx': 'document', '.doc': 'document', '.odt': 'document',
        '.pptx': 'presentation', '.ppt': 'presentation', '.odp': 'presentation',
        '.xlsx': 'spreadsheet', '.xls': 'spreadsheet', '.ods': 'spreadsheet',
        '.pdf': 'pdf',
    }

    format_type = doc_formats.get(suffix, 'code')
    is_binary = suffix in doc_formats

    # Read or convert file content
    try:
        if format_type == 'document':
            # Convert Word/ODT documents
            content = _convert_document(actual_path)
            lines = content.split('\n')
        elif format_type == 'presentation':
            # Convert PowerPoint/ODP presentations
            content = _convert_presentation(actual_path)
            lines = content.split('\n')
        elif format_type == 'spreadsheet':
            # Convert Excel/ODS spreadsheets
            content = _convert_spreadsheet(actual_path)
            lines = content.split('\n')
        elif format_type == 'pdf':
            # Convert PDF (if pdftotext available)
            content = _convert_pdf(actual_path)
            lines = content.split('\n')
        else:
            # Plain text/code files
            content = actual_path.read_text(encoding='utf-8', errors='replace')
            lines = content.split('\n')
    except Exception as e:
        logger.error(f"Failed to read/convert file: {e}")
        raise HTTPException(500, f"Failed to read file: {e}")

    # Get chunks for this file from Project RAG
    chunks = []
    try:
        project = _get_rag_project(path)
        if project and project.is_initialized():
            rel_path = str(actual_path.relative_to(project_root))

            # Method 1: Use metadata store (most reliable)
            meta_store = project.metadata_store
            if meta_store:
                # Get file metadata first
                file_meta = meta_store.get_file_by_path(rel_path)
                if file_meta:
                    # Load all chunks and filter by file_id
                    all_chunks = meta_store.load_chunks()
                    for chunk_id, chunk_meta in all_chunks.items():
                        if chunk_meta.file_id == file_meta.file_id:
                            chunks.append({
                                'chunk_id': chunk_meta.chunk_id,
                                'line_start': chunk_meta.line_start,
                                'line_end': chunk_meta.line_end,
                                'content': chunk_meta.text_preview[:200] if chunk_meta.text_preview else '',
                                'concepts': chunk_meta.tags,
                                'score': 1.0,
                            })

            # Method 2: Fallback to vector store query if metadata didn't work
            if not chunks:
                vs = project.vector_store
                if vs:
                    # Query with file path filter
                    results = vs.query_mixed(
                        query_text=rel_path,
                        top_k=100,
                        min_score=0.0
                    )
                    for r in results:
                        meta = r.get('metadata', {})
                        chunk_file = meta.get('file_path', '')
                        if chunk_file == rel_path or chunk_file.endswith(file_p.name):
                            chunks.append({
                                'chunk_id': r.get('id', ''),
                                'line_start': meta.get('line_start', 0),
                                'line_end': meta.get('line_end', 0),
                                'content': r.get('content', '')[:200],
                                'concepts': meta.get('concepts', []),
                                'score': r.get('score', 0),
                            })

            # Sort chunks by line_start for proper ordering
            chunks.sort(key=lambda c: (c.get('line_start', 0), c.get('line_end', 0)))

            # Also get chunk-concept associations if concept specified
            if concept:
                graph = project.graph
                if graph:
                    # Find chunks that mention this concept
                    concept_chunks = graph.get_concept_chunks(concept) if hasattr(graph, 'get_concept_chunks') else []
                    for chunk in chunks:
                        chunk['highlight'] = any(c.get('chunk_id') == chunk['chunk_id'] for c in concept_chunks)
    except Exception as e:
        logger.warning(f"Could not get chunks for file: {e}")

    # Determine file type for syntax highlighting
    suffix = actual_path.suffix.lower()
    lang_map = {
        '.py': 'python', '.js': 'javascript', '.ts': 'typescript',
        '.java': 'java', '.c': 'c', '.cpp': 'cpp', '.h': 'c',
        '.html': 'html', '.css': 'css', '.json': 'json',
        '.xml': 'xml', '.yaml': 'yaml', '.yml': 'yaml',
        '.md': 'markdown', '.sh': 'bash', '.sql': 'sql',
    }
    language = lang_map.get(suffix, 'plaintext')

    # Extended language map for documents
    doc_lang_map = {
        '.docx': 'markdown', '.doc': 'markdown', '.odt': 'markdown',
        '.pptx': 'markdown', '.ppt': 'markdown', '.odp': 'markdown',
        '.xlsx': 'markdown', '.xls': 'markdown', '.ods': 'markdown',
        '.pdf': 'plaintext',
    }
    if suffix in doc_lang_map:
        language = doc_lang_map[suffix]

    return {
        'file_path': str(actual_path.relative_to(project_root)),
        'absolute_path': str(actual_path),
        'file_name': actual_path.name,
        'language': language,
        'format_type': format_type,  # 'code', 'document', 'presentation', 'spreadsheet', 'pdf'
        'is_binary': is_binary,
        'line_count': len(lines),
        'content': content,
        'lines': lines,
        'chunks': sorted(chunks, key=lambda c: c.get('line_start', 0)),
        'concept': concept,
    }


class ConceptCreateRequest(BaseModel):
    """Request to create a new concept."""
    label: str
    description: str = ""
    origin: str = "user"  # "user" or "auto"


@router.post("/concept/create")
async def create_concept(
    request: ConceptCreateRequest,
    project_path: Optional[str] = Query(None, description="Project path (uses current if not provided)"),
) -> Dict[str, Any]:
    """
    Create or get a concept node in the knowledge graph.

    If the concept already exists, returns the existing node.
    """
    path = project_path or _current_project_path
    project = _get_rag_project(path)
    if not project:
        raise HTTPException(503, "Project RAG module not available or no project selected")

    if not project.is_initialized():
        return {"error": "Project RAG not initialized"}

    try:
        graph = project.graph
        node = graph.get_or_create_concept(
            label=request.label,
            origin=request.origin,
            description=request.description,
        )

        # Save graph to persist the new concept
        graph.save()

        return {
            "status": "created" if node else "exists",
            "concept": {
                "id": node.id if node else None,
                "label": request.label,
                "origin": request.origin,
                "description": request.description,
            },
        }

    except Exception as e:
        logger.error(f"Concept creation failed: {e}")
        raise HTTPException(500, f"Concept creation failed: {e}")


# =============================================================================
# Chat Context Integration
# =============================================================================

def retrieve_project_rag_context(
    query: str,
    top_k: int = 5,
    max_chars_per_chunk: int = 1500,
    collection: str = "mixed",
) -> Optional[str]:
    """
    Retrieve context from Project RAG for chat injection.

    This is called by the chat system to get relevant project context
    before sending the message to the LLM.

    Args:
        query: User query
        top_k: Number of chunks to retrieve
        max_chars_per_chunk: Maximum chars per chunk
        collection: Collection to search (docs, code, mixed)

    Returns:
        Formatted context string or None if not available
    """
    # Log which project is being used for debugging
    logger.info(f"Project RAG retrieval: current_project_path={_current_project_path}")

    project = _get_rag_project()
    if not project:
        logger.info("Project RAG: no project instance available")
        return None

    logger.info(f"Project RAG: using project_root={project.project_root}")

    if not project.is_initialized():
        logger.info(f"Project RAG: not initialized at {project.project_root}")
        return None

    try:
        context = project.retrieve_context(
            query=query,
            top_k=top_k,
            collection=collection,
            max_chars_per_chunk=max_chars_per_chunk,
        )

        if context.is_empty():
            logger.info("Project RAG: no results found")
            return None

        # Log retrieved files for debugging
        for i, result in enumerate(context.results):
            logger.info(f"Project RAG result {i+1}: {result.file_path}")

        return context.format_for_prompt()

    except Exception as e:
        logger.warning(f"Project RAG retrieval failed: {e}")
        return None


def check_project_rag_available() -> bool:
    """Check if Project RAG is available and initialized."""
    project = _get_rag_project()
    return project is not None and project.is_initialized()
