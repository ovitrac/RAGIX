"""
RAGIX RAG Router - RAG (Retrieval-Augmented Generation) management endpoints

Provides API for managing RAG indexes, viewing indexed documents, and controlling
RAG activation.

Author: Olivier Vitrac, PhD, HDR | olivier.vitrac@adservio.fr | Adservio | 2025-12-05
"""

import io
import json
import os
import shutil
import subprocess
import tempfile
import zipfile
from datetime import datetime
from pathlib import Path
from typing import Dict, Any, List, Optional

from fastapi import APIRouter, HTTPException, Query, UploadFile, File, Form
from pydantic import BaseModel

router = APIRouter(prefix="/api/rag", tags=["rag"])

# RAG state (per-session)
_rag_state: Dict[str, Dict[str, Any]] = {}
_storage_root: str = ""


def set_rag_store(storage_root: str = ""):
    """Set the storage root from server.py."""
    global _storage_root
    _storage_root = storage_root


def _get_rag_state(session_id: str) -> Dict[str, Any]:
    """Get or initialize RAG state for a session."""
    if session_id not in _rag_state:
        _rag_state[session_id] = {
            "enabled": False,
            "index_loaded": False,
            "index_path": None,
            "document_count": 0,
            "chunk_count": 0,
            "last_indexed": None,
            "embedding_model": None,
            "fusion_strategy": "rrf",
        }
    return _rag_state[session_id]


def _get_index_path() -> Path:
    """Get the RAG index path."""
    root = Path(_storage_root) if _storage_root else Path.cwd()
    return root / ".ragix" / "index"


# Request/Response Models

class RAGEnableRequest(BaseModel):
    """Request to enable/disable RAG."""
    enabled: bool


class RAGConfigRequest(BaseModel):
    """Request to update RAG configuration."""
    fusion_strategy: Optional[str] = None
    top_k: Optional[int] = None
    bm25_weight: Optional[float] = None
    vector_weight: Optional[float] = None


class RAGIndexRequest(BaseModel):
    """Request to index documents."""
    paths: List[str]
    recursive: bool = True
    extensions: Optional[List[str]] = None  # e.g., [".py", ".md", ".txt"]


# Endpoints

@router.get("/status")
async def get_rag_status(session_id: str = Query("default")) -> Dict[str, Any]:
    """
    Get RAG status for a session.

    Returns:
        enabled: Whether RAG is enabled
        index_loaded: Whether an index is loaded
        document_count: Number of indexed documents
        chunk_count: Number of indexed chunks
        index_path: Path to the index directory
        fusion_strategy: Current fusion strategy
    """
    state = _get_rag_state(session_id)
    index_path = _get_index_path()

    # Check if index exists on disk
    bm25_exists = (index_path / "bm25_index.json").exists()
    vector_exists = (index_path / "vector_index.npz").exists() or (index_path / "embeddings.pkl").exists()

    return {
        "session_id": session_id,
        "enabled": state["enabled"],
        "index_loaded": state["index_loaded"],
        "index_exists": bm25_exists or vector_exists,
        "bm25_index_exists": bm25_exists,
        "vector_index_exists": vector_exists,
        "document_count": state["document_count"],
        "chunk_count": state["chunk_count"],
        "last_indexed": state["last_indexed"],
        "index_path": str(index_path),
        "fusion_strategy": state["fusion_strategy"],
        "embedding_model": state["embedding_model"],
    }


@router.post("/enable")
async def set_rag_enabled(
    request: RAGEnableRequest,
    session_id: str = Query("default"),
) -> Dict[str, Any]:
    """Enable or disable RAG for a session."""
    state = _get_rag_state(session_id)
    state["enabled"] = request.enabled

    if request.enabled and not state["index_loaded"]:
        # Try to load existing index
        index_path = _get_index_path()
        if (index_path / "bm25_index.json").exists():
            # Index exists, mark as loaded (actual loading done lazily)
            state["index_loaded"] = True
            # TODO: Load index and count documents/chunks

    return {
        "session_id": session_id,
        "enabled": state["enabled"],
        "index_loaded": state["index_loaded"],
        "status": "enabled" if state["enabled"] else "disabled",
    }


@router.get("/config")
async def get_rag_config(session_id: str = Query("default")) -> Dict[str, Any]:
    """Get RAG configuration."""
    state = _get_rag_state(session_id)

    # Load from ragix.yaml if available
    try:
        from ragix_core.config import get_config
        config = get_config()
        search_config = config.search if hasattr(config, 'search') and config.search else None

        return {
            "session_id": session_id,
            "enabled": state["enabled"],
            "fusion_strategy": search_config.fusion_strategy if search_config else "rrf",
            "top_k": search_config.top_k if search_config else 10,
            "bm25_weight": search_config.bm25_weight if search_config else 0.5,
            "vector_weight": search_config.vector_weight if search_config else 0.5,
            "embedding_model": search_config.embedding_model if search_config else "all-MiniLM-L6-v2",
            "index_path": str(_get_index_path()),
        }
    except Exception as e:
        return {
            "session_id": session_id,
            "enabled": state["enabled"],
            "fusion_strategy": state["fusion_strategy"],
            "error": str(e),
        }


@router.post("/config")
async def update_rag_config(
    request: RAGConfigRequest,
    session_id: str = Query("default"),
) -> Dict[str, Any]:
    """Update RAG configuration (session-level, does not persist to ragix.yaml)."""
    state = _get_rag_state(session_id)

    if request.fusion_strategy:
        valid_strategies = ["rrf", "weighted", "interleave", "bm25_only", "vector_only"]
        if request.fusion_strategy not in valid_strategies:
            raise HTTPException(400, f"Invalid fusion strategy. Valid: {valid_strategies}")
        state["fusion_strategy"] = request.fusion_strategy

    return {
        "session_id": session_id,
        "fusion_strategy": state["fusion_strategy"],
        "status": "updated",
    }


@router.get("/documents")
async def list_indexed_documents(
    session_id: str = Query("default"),
    limit: int = Query(50, ge=1, le=500),
    offset: int = Query(0, ge=0),
) -> Dict[str, Any]:
    """
    List indexed documents.

    Returns summary of documents in the RAG index.
    """
    index_path = _get_index_path()
    metadata_path = index_path / "metadata.json"

    if not metadata_path.exists():
        return {
            "session_id": session_id,
            "documents": [],
            "total": 0,
            "message": "No index found. Use /api/rag/index to create one.",
        }

    try:
        import json
        with open(metadata_path, 'r') as f:
            metadata = json.load(f)

        documents = metadata.get("documents", [])
        total = len(documents)

        # Apply pagination
        paginated = documents[offset:offset + limit]

        return {
            "session_id": session_id,
            "documents": paginated,
            "total": total,
            "offset": offset,
            "limit": limit,
            "indexed_at": metadata.get("indexed_at"),
        }
    except Exception as e:
        return {
            "session_id": session_id,
            "documents": [],
            "total": 0,
            "error": str(e),
        }


@router.get("/chunks")
async def list_indexed_chunks(
    session_id: str = Query("default"),
    file_path: Optional[str] = None,
    limit: int = Query(50, ge=1, le=200),
    offset: int = Query(0, ge=0),
) -> Dict[str, Any]:
    """
    List indexed chunks, optionally filtered by file path.

    Returns chunk details from the RAG index.
    """
    index_path = _get_index_path()
    chunks_path = index_path / "chunks.json"

    if not chunks_path.exists():
        return {
            "session_id": session_id,
            "chunks": [],
            "total": 0,
            "message": "No chunks found. Use /api/rag/index to create an index.",
        }

    try:
        import json
        with open(chunks_path, 'r') as f:
            all_chunks = json.load(f)

        # Filter by file path if specified
        if file_path:
            chunks = [c for c in all_chunks if c.get("file_path", "").endswith(file_path)]
        else:
            chunks = all_chunks

        total = len(chunks)
        paginated = chunks[offset:offset + limit]

        # Truncate content for display
        for chunk in paginated:
            content = chunk.get("content", "")
            if len(content) > 200:
                chunk["content"] = content[:200] + "..."
                chunk["content_truncated"] = True

        return {
            "session_id": session_id,
            "chunks": paginated,
            "total": total,
            "offset": offset,
            "limit": limit,
            "filter_file": file_path,
        }
    except Exception as e:
        return {
            "session_id": session_id,
            "chunks": [],
            "total": 0,
            "error": str(e),
        }


@router.post("/search")
async def search_rag(
    query: str,
    session_id: str = Query("default"),
    top_k: int = Query(10, ge=1, le=50),
) -> Dict[str, Any]:
    """
    Search the RAG index.

    Returns ranked results from hybrid search.
    """
    state = _get_rag_state(session_id)

    if not state["enabled"]:
        return {
            "session_id": session_id,
            "results": [],
            "error": "RAG is not enabled. Use /api/rag/enable to activate.",
        }

    # TODO: Implement actual search using HybridSearchEngine
    # For now, return placeholder
    return {
        "session_id": session_id,
        "query": query,
        "top_k": top_k,
        "results": [],
        "message": "Search not yet implemented - index loading required",
    }


@router.delete("/index")
async def clear_rag_index(session_id: str = Query("default")) -> Dict[str, Any]:
    """
    Clear the RAG index.

    Removes all indexed documents and chunks.
    """
    index_path = _get_index_path()

    if not index_path.exists():
        return {
            "session_id": session_id,
            "status": "no_index",
            "message": "No index to clear.",
        }

    import shutil
    try:
        shutil.rmtree(index_path)
        state = _get_rag_state(session_id)
        state["index_loaded"] = False
        state["document_count"] = 0
        state["chunk_count"] = 0

        return {
            "session_id": session_id,
            "status": "cleared",
            "message": f"Index cleared: {index_path}",
        }
    except Exception as e:
        raise HTTPException(500, f"Failed to clear index: {e}")


@router.get("/stats")
async def get_rag_stats(session_id: str = Query("default")) -> Dict[str, Any]:
    """
    Get RAG index statistics.

    Returns detailed statistics about the index.
    """
    index_path = _get_index_path()

    stats = {
        "session_id": session_id,
        "index_path": str(index_path),
        "exists": index_path.exists(),
        "files": {},
        "total_size_bytes": 0,
    }

    if index_path.exists():
        for file in index_path.iterdir():
            if file.is_file():
                size = file.stat().st_size
                stats["files"][file.name] = {
                    "size_bytes": size,
                    "size_kb": round(size / 1024, 2),
                }
                stats["total_size_bytes"] += size

        stats["total_size_kb"] = round(stats["total_size_bytes"] / 1024, 2)

    return stats


# ==============================================================================
# v0.33: RAG Feed Interface - Upload and Index
# ==============================================================================

# Supported text file extensions for indexing
TEXT_EXTENSIONS = {
    # General text
    '.txt', '.md', '.rst', '.log', '.csv',
    # Config files
    '.json', '.yaml', '.yml', '.xml', '.toml', '.ini', '.cfg', '.conf', '.env', '.properties',
    # Web
    '.html', '.htm', '.css', '.scss', '.less',
    # JavaScript/TypeScript
    '.js', '.ts', '.jsx', '.tsx', '.mjs', '.cjs', '.vue', '.svelte',
    # Python
    '.py', '.pyw', '.pyx', '.pxd', '.pyi',
    # Java/JVM
    '.java', '.kt', '.kts', '.groovy', '.gradle', '.scala',
    # C/C++
    '.c', '.h', '.cpp', '.hpp', '.cc', '.hh', '.cxx', '.hxx',
    # Shell/Scripts
    '.sh', '.bash', '.zsh', '.fish', '.ps1', '.bat', '.cmd',
    # MATLAB/Octave
    '.m',
    # Build/Project files
    '.pom', '.makefile', '.cmake', '.dockerfile',
    # SQL
    '.sql', '.ddl', '.dml',
    # Other languages
    '.rb', '.php', '.go', '.rs', '.swift', '.r', '.jl', '.lua', '.pl', '.pm',
    # Git
    '.gitignore', '.gitattributes', '.gitmodules',
}

# v0.33: Document extensions that require conversion to text/markdown
CONVERSION_EXTENSIONS = {
    '.pdf',    # Convert via pdftotext
    '.docx',   # Convert via pandoc
    '.doc',    # Convert via pandoc (may need libreoffice)
    '.odt',    # Convert via pandoc
    '.rtf',    # Convert via pandoc
    '.pptx',   # Convert via python-pptx
    '.ppt',    # Legacy PowerPoint (limited support)
    '.xlsx',   # Convert via openpyxl
    '.xls',    # Legacy Excel (limited support)
}

# Extensions that pandoc can handle directly
PANDOC_EXTENSIONS = {'.docx', '.doc', '.odt', '.rtf', '.epub', '.html'}

# Extensions needing special Python libraries
PPTX_EXTENSIONS = {'.pptx', '.ppt'}
XLSX_EXTENSIONS = {'.xlsx', '.xls'}

# Default RAG indexing parameters
DEFAULT_CHUNK_SIZE = 1000  # characters
DEFAULT_CHUNK_OVERLAP = 200  # characters
DEFAULT_MAX_FILE_SIZE = 10 * 1024 * 1024  # 10 MB


class RAGIndexConfig(BaseModel):
    """Configuration for RAG indexing."""
    chunk_size: int = DEFAULT_CHUNK_SIZE
    chunk_overlap: int = DEFAULT_CHUNK_OVERLAP
    max_file_size: int = DEFAULT_MAX_FILE_SIZE
    extensions: Optional[List[str]] = None  # Custom extension filter


def _chunk_text(text: str, chunk_size: int = DEFAULT_CHUNK_SIZE, overlap: int = DEFAULT_CHUNK_OVERLAP) -> List[Dict[str, Any]]:
    """
    Split text into overlapping chunks.

    Args:
        text: Text to split
        chunk_size: Size of each chunk in characters
        overlap: Overlap between chunks

    Returns:
        List of chunk dictionaries with content and metadata
    """
    chunks = []
    if not text:
        return chunks

    # Split by paragraphs first, then by sentences if needed
    paragraphs = text.split('\n\n')
    current_chunk = ""
    chunk_index = 0

    for para in paragraphs:
        if len(current_chunk) + len(para) <= chunk_size:
            current_chunk += para + "\n\n"
        else:
            # Save current chunk if not empty
            if current_chunk.strip():
                chunks.append({
                    "content": current_chunk.strip(),
                    "chunk_index": chunk_index,
                    "char_start": max(0, len(text) - len(current_chunk) - sum(len(c["content"]) for c in chunks)),
                })
                chunk_index += 1

            # Start new chunk with overlap
            if overlap > 0 and current_chunk:
                # Take last 'overlap' characters from previous chunk
                overlap_text = current_chunk[-overlap:] if len(current_chunk) > overlap else current_chunk
                current_chunk = overlap_text + para + "\n\n"
            else:
                current_chunk = para + "\n\n"

            # Handle very long paragraphs
            while len(current_chunk) > chunk_size:
                chunks.append({
                    "content": current_chunk[:chunk_size].strip(),
                    "chunk_index": chunk_index,
                })
                chunk_index += 1
                current_chunk = current_chunk[chunk_size - overlap:]

    # Add remaining chunk
    if current_chunk.strip():
        chunks.append({
            "content": current_chunk.strip(),
            "chunk_index": chunk_index,
        })

    return chunks


def _is_text_file(filename: str, allowed_extensions: Optional[set] = None) -> bool:
    """Check if file is a text file based on extension."""
    ext = Path(filename).suffix.lower()
    if allowed_extensions:
        return ext in allowed_extensions
    return ext in TEXT_EXTENSIONS


def _needs_conversion(filename: str, pdf_enabled: bool = True, pandoc_enabled: bool = True) -> bool:
    """
    Check if file needs conversion (docx, pdf, pptx, xlsx, etc.).

    Args:
        filename: Filename to check
        pdf_enabled: Whether PDF conversion is enabled
        pandoc_enabled: Whether pandoc (Office) conversion is enabled

    Returns:
        True if file needs conversion and converter is enabled
    """
    ext = Path(filename).suffix.lower()
    if ext not in CONVERSION_EXTENSIONS:
        return False

    # Check if the specific converter is enabled
    if ext == '.pdf':
        return pdf_enabled
    else:
        # Office formats: docx, doc, odt, rtf, pptx, xlsx
        return pandoc_enabled


def _convert_document_to_text(file_path: str, ext: str) -> str:
    """
    Convert a document file to plain text or markdown.

    Uses pdftotext for PDF files, pandoc for Office documents,
    and Python libraries for PPTX/XLSX.

    Args:
        file_path: Path to the file to convert
        ext: File extension (lowercase, with dot)

    Returns:
        Extracted text content

    Raises:
        Exception: If conversion fails
    """
    try:
        from ragix_core.config import get_config
        config = get_config()
        conv_config = config.converters
    except Exception:
        # Default configuration if ragix_core not available
        class DefaultConfig:
            timeout = 60
            class pdftotext:
                enabled = True
                path = 'pdftotext'
                options = ['-layout']
            class pandoc:
                enabled = True
                path = 'pandoc'
                output_format = 'markdown'  # Use markdown for better RAG
                options = []
        conv_config = DefaultConfig()

    extracted_text = ""

    if ext == '.pdf':
        # Use pdftotext for PDF conversion
        pdftotext_path = getattr(conv_config.pdftotext, 'path', 'pdftotext') or 'pdftotext'
        if not shutil.which(pdftotext_path):
            raise Exception(f"pdftotext not found. Install with: sudo apt install poppler-utils")

        options = getattr(conv_config.pdftotext, 'options', ['-layout']) or ['-layout']
        cmd = [pdftotext_path] + options + [file_path, '-']

        result = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            timeout=getattr(conv_config, 'timeout', 60)
        )

        if result.returncode != 0:
            raise Exception(f"PDF conversion failed: {result.stderr}")

        extracted_text = result.stdout

    elif ext in PPTX_EXTENSIONS:
        # Use python-pptx for PowerPoint files
        extracted_text = _convert_pptx_to_text(file_path)

    elif ext in XLSX_EXTENSIONS:
        # Use openpyxl for Excel files
        extracted_text = _convert_xlsx_to_text(file_path)

    elif ext in PANDOC_EXTENSIONS:
        # Use pandoc for docx, doc, odt, rtf
        pandoc_path = getattr(conv_config.pandoc, 'path', 'pandoc') or 'pandoc'
        if not shutil.which(pandoc_path):
            raise Exception(f"pandoc not found. Install with: sudo apt install pandoc")

        # Use markdown output for better RAG indexing
        output_format = getattr(conv_config.pandoc, 'output_format', 'markdown') or 'markdown'
        options = getattr(conv_config.pandoc, 'options', []) or []
        cmd = [pandoc_path, '-t', output_format] + options + [file_path]

        result = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            timeout=getattr(conv_config, 'timeout', 60)
        )

        if result.returncode != 0:
            raise Exception(f"Document conversion failed: {result.stderr}")

        extracted_text = result.stdout

    else:
        raise Exception(f"Unsupported conversion format: {ext}")

    return extracted_text.strip()


def _convert_pptx_to_text(file_path: str) -> str:
    """
    Convert PowerPoint file to text using python-pptx.

    Args:
        file_path: Path to the PPTX file

    Returns:
        Extracted text content
    """
    try:
        from pptx import Presentation
    except ImportError:
        raise Exception("python-pptx not installed. Install with: pip install python-pptx")

    prs = Presentation(file_path)
    text_parts = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = [f"\n## Slide {slide_num}\n"]

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())

            # Handle tables
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    slide_text.append(row_text)

        if len(slide_text) > 1:  # More than just the header
            text_parts.append("\n".join(slide_text))

    return "\n\n".join(text_parts)


def _convert_xlsx_to_text(file_path: str) -> str:
    """
    Convert Excel file to text using openpyxl.

    Args:
        file_path: Path to the XLSX file

    Returns:
        Extracted text content
    """
    try:
        from openpyxl import load_workbook
    except ImportError:
        raise Exception("openpyxl not installed. Install with: pip install openpyxl")

    wb = load_workbook(file_path, read_only=True, data_only=True)
    text_parts = []

    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        sheet_text = [f"\n## Sheet: {sheet_name}\n"]

        for row in sheet.iter_rows(values_only=True):
            # Filter out empty rows
            row_values = [str(cell) if cell is not None else "" for cell in row]
            if any(v.strip() for v in row_values):
                sheet_text.append(" | ".join(row_values))

        if len(sheet_text) > 1:
            text_parts.append("\n".join(sheet_text))

    wb.close()
    return "\n\n".join(text_parts)


def _convert_bytes_to_text(content: bytes, filename: str) -> str:
    """
    Convert document bytes to text.

    Saves content to a temp file and calls _convert_document_to_text.

    Args:
        content: File content as bytes
        filename: Original filename (for extension detection)

    Returns:
        Extracted text content
    """
    ext = Path(filename).suffix.lower()

    # Save to temp file
    with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp_file:
        tmp_file.write(content)
        tmp_path = tmp_file.name

    try:
        return _convert_document_to_text(tmp_path, ext)
    finally:
        try:
            os.unlink(tmp_path)
        except:
            pass


def _build_bm25_index(chunks: List[Dict[str, Any]], index_path: Path) -> bool:
    """
    Build BM25 index from chunks for actual search capability.

    Args:
        chunks: List of chunk dictionaries with 'content' and 'file_path'
        index_path: Path to save the BM25 index

    Returns:
        True if index was built successfully
    """
    try:
        from ragix_core.bm25_index import BM25Index, BM25Document
    except ImportError:
        # Fallback: create a simple JSON-based index
        return _build_simple_bm25_index(chunks, index_path)

    # Create BM25 index
    bm25 = BM25Index()

    for i, chunk in enumerate(chunks):
        content = chunk.get("content", "")
        file_path = chunk.get("file_path", f"chunk_{i}")
        chunk_idx = chunk.get("chunk_index", i)

        # Tokenize content
        tokens = bm25.tokenizer.tokenize(content)

        doc = BM25Document(
            doc_id=f"{file_path}:{chunk_idx}",
            file_path=file_path,
            start_line=chunk_idx * 10,  # Approximate
            end_line=(chunk_idx + 1) * 10,
            chunk_type="text",
            name=f"chunk_{chunk_idx}",
            tokens=tokens,
        )
        bm25.add_document(doc)

    # Save the index
    bm25.save(index_path)
    return True


def _build_simple_bm25_index(chunks: List[Dict[str, Any]], index_path: Path) -> bool:
    """
    Build a simple JSON-based BM25-compatible index when ragix_core is not available.

    Args:
        chunks: List of chunk dictionaries
        index_path: Path to save the index

    Returns:
        True if index was built successfully
    """
    import re
    from collections import Counter

    # Simple tokenizer
    def tokenize(text: str) -> List[str]:
        text = text.lower()
        tokens = re.findall(r'\b[a-z_][a-z0-9_]*\b', text)
        return [t for t in tokens if len(t) >= 2]

    docs_data = {}
    for i, chunk in enumerate(chunks):
        content = chunk.get("content", "")
        file_path = chunk.get("file_path", f"chunk_{i}")
        chunk_idx = chunk.get("chunk_index", i)
        doc_id = f"{file_path}:{chunk_idx}"

        tokens = tokenize(content)
        docs_data[doc_id] = {
            "file_path": file_path,
            "start_line": chunk_idx * 10,
            "end_line": (chunk_idx + 1) * 10,
            "chunk_type": "text",
            "name": f"chunk_{chunk_idx}",
            "tokens": tokens,
            "content": content[:500],  # Store truncated content for retrieval
        }

    # Calculate statistics
    total_tokens = sum(len(d["tokens"]) for d in docs_data.values())
    avg_doc_length = total_tokens / len(docs_data) if docs_data else 0

    # Save documents
    with open(index_path / "bm25_documents.json", "w", encoding="utf-8") as f:
        json.dump(docs_data, f)

    # Save metadata
    metadata = {
        "k1": 1.5,
        "b": 0.75,
        "total_docs": len(docs_data),
        "avg_doc_length": avg_doc_length,
        "total_tokens": total_tokens,
        "vocabulary_size": len(set(t for d in docs_data.values() for t in d["tokens"])),
    }

    with open(index_path / "bm25_metadata.json", "w", encoding="utf-8") as f:
        json.dump(metadata, f, indent=2)

    return True


@router.post("/upload")
async def upload_and_index(
    file: UploadFile = File(...),
    session_id: str = Form("default"),
    chunk_size: int = Form(DEFAULT_CHUNK_SIZE),
    chunk_overlap: int = Form(DEFAULT_CHUNK_OVERLAP),
    pdf_enabled: str = Form("true"),
    pandoc_enabled: str = Form("true"),
) -> Dict[str, Any]:
    """
    Upload a file or archive and index it for RAG.

    Supports:
    - Single text files
    - ZIP archives containing multiple files
    - Document files (PDF, DOCX, PPTX, XLSX) with conversion

    Args:
        file: File to upload (text file or ZIP)
        session_id: Session ID
        chunk_size: Size of text chunks (default: 1000 chars)
        chunk_overlap: Overlap between chunks (default: 200 chars)
        pdf_enabled: Enable PDF conversion via pdftotext
        pandoc_enabled: Enable Office docs conversion via pandoc

    Returns:
        Indexing statistics
    """
    # Parse boolean strings from form data
    pdf_conv_enabled = pdf_enabled.lower() in ('true', '1', 'yes', 'on')
    pandoc_conv_enabled = pandoc_enabled.lower() in ('true', '1', 'yes', 'on')

    index_path = _get_index_path()
    index_path.mkdir(parents=True, exist_ok=True)

    files_indexed = 0
    files_skipped = 0
    chunks_created = 0
    errors = []

    # Load or create metadata
    metadata_path = index_path / "metadata.json"
    chunks_path = index_path / "chunks.json"

    if metadata_path.exists():
        with open(metadata_path, 'r') as f:
            metadata = json.load(f)
    else:
        metadata = {"documents": [], "indexed_at": None, "config": {}}

    if chunks_path.exists():
        with open(chunks_path, 'r') as f:
            all_chunks = json.load(f)
    else:
        all_chunks = []

    # Determine file type
    filename = file.filename or "upload"
    content = await file.read()

    files_converted = 0  # Track documents converted

    if filename.lower().endswith('.zip'):
        # Handle ZIP archive
        try:
            with zipfile.ZipFile(io.BytesIO(content)) as zf:
                for zip_info in zf.infolist():
                    if zip_info.is_dir():
                        continue

                    inner_filename = zip_info.filename
                    is_text = _is_text_file(inner_filename)
                    needs_conv = _needs_conversion(inner_filename, pdf_conv_enabled, pandoc_conv_enabled)

                    if not is_text and not needs_conv:
                        files_skipped += 1
                        continue

                    if zip_info.file_size > DEFAULT_MAX_FILE_SIZE:
                        errors.append(f"File too large: {inner_filename}")
                        files_skipped += 1
                        continue

                    try:
                        raw_content = zf.read(inner_filename)

                        # v0.33: Convert document files (docx, pdf, pptx, xlsx) to text
                        if needs_conv:
                            try:
                                file_content = _convert_bytes_to_text(raw_content, inner_filename)
                                files_converted += 1
                            except Exception as conv_err:
                                errors.append(f"Conversion failed for {inner_filename}: {str(conv_err)}")
                                files_skipped += 1
                                continue
                        else:
                            file_content = raw_content.decode('utf-8', errors='replace')

                        chunks = _chunk_text(file_content, chunk_size, chunk_overlap)

                        for chunk in chunks:
                            chunk["file_path"] = inner_filename
                            chunk["source_archive"] = filename
                            chunk["converted"] = needs_conv
                            chunk["indexed_at"] = datetime.now().isoformat()
                            all_chunks.append(chunk)
                            chunks_created += 1

                        metadata["documents"].append({
                            "path": inner_filename,
                            "source": filename,
                            "size": zip_info.file_size,
                            "chunks": len(chunks),
                            "converted": needs_conv,
                            "indexed_at": datetime.now().isoformat(),
                        })
                        files_indexed += 1

                    except Exception as e:
                        errors.append(f"Error reading {inner_filename}: {str(e)}")
                        files_skipped += 1

        except zipfile.BadZipFile:
            raise HTTPException(400, "Invalid ZIP file")

    else:
        # Handle single file
        is_text = _is_text_file(filename)
        needs_conv = _needs_conversion(filename, pdf_conv_enabled, pandoc_conv_enabled)

        if not is_text and not needs_conv:
            raise HTTPException(400, f"Unsupported file type: {filename}")

        try:
            # v0.33: Convert document files (docx, pdf, pptx, xlsx) to text/markdown
            if needs_conv:
                try:
                    file_content = _convert_bytes_to_text(content, filename)
                    files_converted += 1
                except Exception as conv_err:
                    raise HTTPException(500, f"Conversion failed: {str(conv_err)}")
            else:
                file_content = content.decode('utf-8', errors='replace')

            chunks = _chunk_text(file_content, chunk_size, chunk_overlap)

            for chunk in chunks:
                chunk["file_path"] = filename
                chunk["converted"] = needs_conv
                chunk["indexed_at"] = datetime.now().isoformat()
                all_chunks.append(chunk)
                chunks_created += 1

            metadata["documents"].append({
                "path": filename,
                "size": len(content),
                "chunks": len(chunks),
                "converted": needs_conv,
                "indexed_at": datetime.now().isoformat(),
            })
            files_indexed += 1

        except HTTPException:
            raise
        except Exception as e:
            raise HTTPException(500, f"Error processing file: {str(e)}")

    # Save updated index
    metadata["indexed_at"] = datetime.now().isoformat()
    metadata["config"] = {
        "chunk_size": chunk_size,
        "chunk_overlap": chunk_overlap,
    }

    with open(metadata_path, 'w') as f:
        json.dump(metadata, f, indent=2)

    with open(chunks_path, 'w') as f:
        json.dump(all_chunks, f, indent=2)

    # v0.33: Build BM25 index from chunks for actual search capability
    bm25_built = False
    try:
        bm25_built = _build_bm25_index(all_chunks, index_path)
    except Exception as e:
        errors.append(f"BM25 index build failed: {str(e)}")

    # Update RAG state
    state = _get_rag_state(session_id)
    state["document_count"] = len(metadata["documents"])
    state["chunk_count"] = len(all_chunks)
    state["last_indexed"] = metadata["indexed_at"]
    state["index_loaded"] = bm25_built

    return {
        "session_id": session_id,
        "status": "indexed",
        "files_indexed": files_indexed,
        "files_converted": files_converted,
        "files_skipped": files_skipped,
        "chunks_created": chunks_created,
        "total_documents": len(metadata["documents"]),
        "total_chunks": len(all_chunks),
        "bm25_built": bm25_built,
        "errors": errors if errors else None,
    }


@router.post("/index-chat")
async def index_chat_history(
    session_id: str = Query("default"),
    chunk_size: int = Query(DEFAULT_CHUNK_SIZE),
) -> Dict[str, Any]:
    """
    Index the current chat history as RAG documents.

    This allows the agent to reference past conversations.

    Args:
        session_id: Session ID
        chunk_size: Size of text chunks

    Returns:
        Indexing statistics
    """
    # Import here to avoid circular imports
    from .sessions import get_sessions_store

    sessions = get_sessions_store()
    if session_id not in sessions:
        raise HTTPException(404, "Session not found")

    session = sessions[session_id]
    history = session.get("message_history", [])

    if not history:
        return {
            "session_id": session_id,
            "status": "empty",
            "message": "No chat history to index",
        }

    index_path = _get_index_path()
    index_path.mkdir(parents=True, exist_ok=True)

    # Load existing index
    metadata_path = index_path / "metadata.json"
    chunks_path = index_path / "chunks.json"

    if metadata_path.exists():
        with open(metadata_path, 'r') as f:
            metadata = json.load(f)
    else:
        metadata = {"documents": [], "indexed_at": None, "config": {}}

    if chunks_path.exists():
        with open(chunks_path, 'r') as f:
            all_chunks = json.load(f)
    else:
        all_chunks = []

    # Convert chat history to text and chunk
    chat_text = ""
    for msg in history:
        role = msg.get("role", "unknown")
        content = msg.get("content", "")
        timestamp = msg.get("timestamp", "")
        chat_text += f"[{role.upper()}] ({timestamp})\n{content}\n\n"

    chunks = _chunk_text(chat_text, chunk_size, chunk_size // 5)  # Less overlap for chat
    chunks_created = 0

    for chunk in chunks:
        chunk["file_path"] = f"chat_history_{session_id}"
        chunk["source_type"] = "chat"
        chunk["session_id"] = session_id
        chunk["indexed_at"] = datetime.now().isoformat()
        all_chunks.append(chunk)
        chunks_created += 1

    # Update metadata
    metadata["documents"].append({
        "path": f"chat_history_{session_id}",
        "source_type": "chat",
        "message_count": len(history),
        "chunks": chunks_created,
        "indexed_at": datetime.now().isoformat(),
    })
    metadata["indexed_at"] = datetime.now().isoformat()

    # Save
    with open(metadata_path, 'w') as f:
        json.dump(metadata, f, indent=2)

    with open(chunks_path, 'w') as f:
        json.dump(all_chunks, f, indent=2)

    # v0.33: Build BM25 index from chunks
    bm25_built = False
    try:
        bm25_built = _build_bm25_index(all_chunks, index_path)
    except Exception:
        pass

    # Update state
    state = _get_rag_state(session_id)
    state["document_count"] = len(metadata["documents"])
    state["chunk_count"] = len(all_chunks)
    state["last_indexed"] = metadata["indexed_at"]
    state["index_loaded"] = bm25_built

    return {
        "session_id": session_id,
        "status": "indexed",
        "messages_indexed": len(history),
        "messages_processed": len(history),
        "chunks_created": chunks_created,
        "total_documents": len(metadata["documents"]),
        "total_chunks": len(all_chunks),
        "bm25_built": bm25_built,
    }


@router.get("/params")
async def get_rag_params(session_id: str = Query("default")) -> Dict[str, Any]:
    """
    Get RAG indexing parameters.

    Returns current chunking configuration and supported file types.
    """
    return {
        "session_id": session_id,
        "chunk_size": DEFAULT_CHUNK_SIZE,
        "chunk_overlap": DEFAULT_CHUNK_OVERLAP,
        "max_file_size": DEFAULT_MAX_FILE_SIZE,
        "supported_extensions": sorted(list(TEXT_EXTENSIONS)),
        "conversion_extensions": sorted(list(CONVERSION_EXTENSIONS)),
        "all_extensions": sorted(list(TEXT_EXTENSIONS | CONVERSION_EXTENSIONS)),
    }
